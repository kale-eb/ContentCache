import os
import json
import tempfile
import subprocess
from pathlib import Path
from datetime import datetime, timezone
from dotenv import load_dotenv
from config import get_text_metadata_path
import api_client

load_dotenv()

class TextProcessor:
    def __init__(self):
        self.supported_extensions = {
            # Direct text processing
            '.txt': self._read_text_file,
            '.md': self._read_text_file,
            '.rtf': self._read_text_file,
            '.csv': self._read_text_file,
            '.json': self._read_text_file,
            '.xml': self._read_text_file,
            '.html': self._read_text_file,
            '.htm': self._read_text_file,
            
            # Document formats requiring conversion
            '.pdf': self._extract_from_pdf,
            '.docx': self._extract_from_docx,
            '.doc': self._extract_from_doc,
            '.pptx': self._extract_from_pptx,
            '.ppt': self._extract_from_ppt,
            '.xlsx': self._extract_from_xlsx,
            '.xls': self._extract_from_xls,
            '.odt': self._extract_from_odt,
            '.ods': self._extract_from_ods,
            '.odp': self._extract_from_odp,
        }

    def _read_text_file(self, file_path):
        """Read plain text files directly."""
        try:
            # Try different encodings
            encodings = ['utf-8', 'utf-16', 'latin-1', 'cp1252']
            
            for encoding in encodings:
                try:
                    with open(file_path, 'r', encoding=encoding) as f:
                        content = f.read()
                    return content
                except UnicodeDecodeError:
                    continue
                    
            # If all encodings fail, read as binary and decode with errors='ignore'
            with open(file_path, 'rb') as f:
                content = f.read().decode('utf-8', errors='ignore')
            return content
            
        except Exception as e:
            print(f"⚠️ Error reading text file: {e}")
            return None

    def _extract_from_pdf(self, file_path):
        """Extract text from PDF using pdfplumber or pypdf."""
        try:
            # Try pdfplumber first (better formatting)
            try:
                import pdfplumber
                text_content = []
                with pdfplumber.open(file_path) as pdf:
                    for page in pdf.pages:
                        text = page.extract_text()
                        if text:
                            text_content.append(text)
                return '\n\n'.join(text_content)
            except ImportError:
                print("⚠️ pdfplumber not available, trying pypdf...")
                
            # Fallback to pypdf
            try:
                import pypdf
                text_content = []
                with open(file_path, 'rb') as f:
                    reader = pypdf.PdfReader(f)
                    for page in reader.pages:
                        text = page.extract_text()
                        if text:
                            text_content.append(text)
                return '\n\n'.join(text_content)
            except ImportError:
                print("⚠️ pypdf not available, trying command line...")
                
            # Command line fallback using pdftotext (poppler-utils)
            try:
                result = subprocess.run([
                    'pdftotext', file_path, '-'
                ], capture_output=True, text=True, timeout=30)
                
                if result.returncode == 0:
                    return result.stdout
                else:
                    print(f"⚠️ pdftotext failed: {result.stderr}")
                    
            except (subprocess.TimeoutExpired, FileNotFoundError):
                print("⚠️ pdftotext command not available")
                
            return None
            
        except Exception as e:
            print(f"⚠️ Error extracting from PDF: {e}")
            return None

    def _extract_from_docx(self, file_path):
        """Extract text from DOCX files."""
        try:
            import docx
            doc = docx.Document(file_path)
            text_content = []
            
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    text_content.append(paragraph.text)
                    
            return '\n\n'.join(text_content)
            
        except ImportError:
            print("⚠️ python-docx not available, trying command line conversion...")
            return self._convert_with_pandoc(file_path, 'docx')
        except Exception as e:
            print(f"⚠️ Error extracting from DOCX: {e}")
            return None

    def _extract_from_doc(self, file_path):
        """Extract text from DOC files using command line tools."""
        try:
            # Try antiword first
            result = subprocess.run([
                'antiword', file_path
            ], capture_output=True, text=True, timeout=30)
            
            if result.returncode == 0:
                return result.stdout
                
        except (subprocess.TimeoutExpired, FileNotFoundError):
            print("⚠️ antiword not available, trying LibreOffice...")
            
        # Try LibreOffice conversion
        return self._convert_with_libreoffice(file_path, 'txt')

    def _extract_from_pptx(self, file_path):
        """Extract text from PPTX files."""
        try:
            from pptx import Presentation
            prs = Presentation(file_path)
            text_content = []
            
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_text = [f"=== Slide {slide_num} ==="]
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text)
                        
                if len(slide_text) > 1:  # More than just the header
                    text_content.append('\n'.join(slide_text))
                    
            return '\n\n'.join(text_content)
            
        except ImportError:
            print("⚠️ python-pptx not available, trying LibreOffice conversion...")
            return self._convert_with_libreoffice(file_path, 'txt')
        except Exception as e:
            print(f"⚠️ Error extracting from PPTX: {e}")
            return None

    def _extract_from_ppt(self, file_path):
        """Extract text from PPT files using LibreOffice."""
        return self._convert_with_libreoffice(file_path, 'txt')

    def _extract_from_xlsx(self, file_path):
        """Extract text from XLSX files."""
        try:
            import openpyxl
            wb = openpyxl.load_workbook(file_path)
            text_content = []
            
            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]
                sheet_text = [f"=== Sheet: {sheet_name} ==="]
                
                for row in sheet.iter_rows(values_only=True):
                    row_text = []
                    for cell in row:
                        if cell is not None:
                            row_text.append(str(cell))
                    if row_text:
                        sheet_text.append('\t'.join(row_text))
                        
                if len(sheet_text) > 1:
                    text_content.append('\n'.join(sheet_text))
                    
            return '\n\n'.join(text_content)
            
        except ImportError:
            print("⚠️ openpyxl not available, trying LibreOffice conversion...")
            return self._convert_with_libreoffice(file_path, 'txt')
        except Exception as e:
            print(f"⚠️ Error extracting from XLSX: {e}")
            return None

    def _extract_from_xls(self, file_path):
        """Extract text from XLS files using LibreOffice."""
        return self._convert_with_libreoffice(file_path, 'txt')

    def _extract_from_odt(self, file_path):
        """Extract text from ODT files."""
        return self._convert_with_libreoffice(file_path, 'txt')

    def _extract_from_ods(self, file_path):
        """Extract text from ODS files."""
        return self._convert_with_libreoffice(file_path, 'txt')

    def _extract_from_odp(self, file_path):
        """Extract text from ODP files."""
        return self._convert_with_libreoffice(file_path, 'txt')

    def _convert_with_libreoffice(self, file_path, output_format):
        """Convert files using LibreOffice command line."""
        try:
            with tempfile.TemporaryDirectory() as temp_dir:
                # Convert to text using LibreOffice
                result = subprocess.run([
                    'libreoffice', '--headless', '--convert-to', output_format,
                    '--outdir', temp_dir, file_path
                ], capture_output=True, text=True, timeout=60)
                
                if result.returncode == 0:
                    # Find the output file
                    base_name = Path(file_path).stem
                    output_file = os.path.join(temp_dir, f"{base_name}.{output_format}")
                    
                    if os.path.exists(output_file):
                        with open(output_file, 'r', encoding='utf-8') as f:
                            return f.read()
                            
                print(f"⚠️ LibreOffice conversion failed: {result.stderr}")
                return None
                
        except (subprocess.TimeoutExpired, FileNotFoundError) as e:
            print(f"⚠️ LibreOffice conversion error: {e}")
            return None

    def _convert_with_pandoc(self, file_path, input_format):
        """Convert files using Pandoc."""
        try:
            result = subprocess.run([
                'pandoc', '-f', input_format, '-t', 'plain', file_path
            ], capture_output=True, text=True, timeout=30)
            
            if result.returncode == 0:
                return result.stdout
            else:
                print(f"⚠️ Pandoc conversion failed: {result.stderr}")
                return None
                
        except (subprocess.TimeoutExpired, FileNotFoundError) as e:
            print(f"⚠️ Pandoc not available: {e}")
            return None

    def _get_file_timestamp(self, file_path):
        """
        Extract timestamp from file metadata.
        
        Args:
            file_path (str): Path to the file
            
        Returns:
            str: ISO format timestamp (preferring creation time, falling back to modification time)
        """
        try:
            import platform
            
            # Get file stats
            stat = os.stat(file_path)
            
            # Try to get creation time (works on Windows and macOS)
            if platform.system() == 'Windows':
                # Windows: use creation time
                timestamp = datetime.fromtimestamp(stat.st_ctime, tz=timezone.utc)
            elif platform.system() == 'Darwin':  # macOS
                # macOS: use birth time if available, otherwise creation time
                if hasattr(stat, 'st_birthtime'):
                    timestamp = datetime.fromtimestamp(stat.st_birthtime, tz=timezone.utc)
                else:
                    timestamp = datetime.fromtimestamp(stat.st_ctime, tz=timezone.utc)
            else:
                # Linux/Unix: use modification time (creation time not reliably available)
                timestamp = datetime.fromtimestamp(stat.st_mtime, tz=timezone.utc)
            
            return timestamp.isoformat()
            
        except Exception as e:
            print(f"⚠️ Could not extract file timestamp: {e}")
            # Fallback to current time if file timestamp extraction fails
            return datetime.now(timezone.utc).isoformat()

    def extract_text_content(self, file_path):
        """
        Extract text content from various file formats.
        
        Args:
            file_path (str): Path to the file
            
        Returns:
            str: Extracted text content or None if failed
        """
        file_ext = Path(file_path).suffix.lower()
        
        if file_ext not in self.supported_extensions:
            print(f"⚠️ Unsupported file type: {file_ext}")
            return None
            
        print(f"📄 Extracting text from {file_ext} file...")
        extractor = self.supported_extensions[file_ext]
        content = extractor(file_path)
        
        if content:
            # Limit content size for API calls
            max_chars = 100000  # ~100KB of text
            if len(content) > max_chars:
                print(f"⚠️ Content too large ({len(content)} chars), truncating to {max_chars} chars...")
                content = content[:max_chars] + "\n\n[CONTENT TRUNCATED...]"
                
        return content

    def analyze_with_gpt4o_mini(self, file_path, text_content):
        """
        Analyze text content using GPT-4o mini via API server with centralized prompts.
        
        Args:
            file_path (str): Path to the file
            text_content (str): Extracted text content
            
        Returns:
            dict: Structured analysis results
        """
        try:
            # Use API server for text analysis with built-in prompts and function schema
            api_response = api_client.get_api_client().openai_text_analysis(file_path, text_content)
            
            if api_response and 'result' in api_response:
                analysis = json.loads(api_response['result'])
                # Add file timestamp from metadata
                analysis["timestamp"] = self._get_file_timestamp(file_path)
                print("✅ [API Server] Text analysis complete")
                return analysis
            else:
                raise Exception("Invalid API response for text analysis")
                
        except Exception as e:
            print(f"❌ API server failed for text analysis: {e}")
            print("⚠️ No fallback available - all prompts are centralized in API server")
            return None

    def process_file(self, file_path):
        """
        Process a single text file completely.
        
        Args:
            file_path (str): Path to text file
            
        Returns:
            dict: Complete analysis results
        """
        print(f"📄 Processing text file: {os.path.basename(file_path)}")
        
        # Step 1: Extract text content
        print("  📖 Extracting text content...")
        text_content = self.extract_text_content(file_path)
        
        if not text_content:
            print("  ❌ Failed to extract text content")
            return None
            
        if not text_content.strip():
            print("  ⚠️ Document appears to be empty")
            return None
            
        print(f"  ✅ Extracted {len(text_content)} characters")
        
        # Step 2: Analyze with GPT-4o mini
        print("  🤖 Analyzing with GPT-4o mini...")
        analysis = self.analyze_with_gpt4o_mini(file_path, text_content)
        
        if not analysis:
            print("  ❌ Failed to get GPT-4o mini analysis")
            return None
            
        # Combine all results
        result = {
            "file_path": os.path.abspath(file_path),
            "filename": os.path.basename(file_path),
            "file_type": Path(file_path).suffix.lower(),
            "file_size": os.path.getsize(file_path),
            "content_length": len(text_content),
            "content_preview": text_content[:500] + "..." if len(text_content) > 500 else text_content,
            "analysis": analysis,
            "processing_timestamp": datetime.now(timezone.utc).isoformat()
        }
        
        # Save to metadata file
        self.save_metadata(result)
        
        # Store result before cleanup to prevent memory leaks
        result_copy = result.copy()
        
        # Use modular cleanup function
        self.cleanup_text_processing(
            text_content=text_content,
            analysis=analysis,
            result=result
        )
        
        print(f"  ✅ Text processing complete")
        
        return result_copy

    def save_metadata(self, metadata):
        """Save metadata to JSON file using config path"""
        output_file = get_text_metadata_path()
        
        # Load existing metadata
        if os.path.exists(output_file):
            with open(output_file, 'r', encoding='utf-8') as f:
                all_metadata = json.load(f)
        else:
            all_metadata = {}
        
        # Add new metadata
        all_metadata[metadata['file_path']] = metadata
        
        # Save updated metadata
        with open(output_file, 'w', encoding='utf-8') as f:
            json.dump(all_metadata, f, indent=2, ensure_ascii=False)
        
        print(f"✅ Text metadata saved to: {output_file}")

    def cleanup_text_processing(self, text_content=None, analysis=None, result=None):
        """
        Cleanup function for text processing to prevent memory leaks.
        
        Args:
            text_content: Extracted text content string
            analysis: GPT-4o mini analysis dict
            result: Final result dict
            
        Returns:
            float: Amount of memory freed in MB
        """
        try:
            import psutil
            import gc
            
            memory_before = psutil.Process().memory_info().rss / 1024 / 1024
            
            # Delete provided variables
            cleanup_vars = {
                'text_content': text_content,
                'analysis': analysis,
                'result': result
            }
            
            deleted_vars = []
            for var_name, var_value in cleanup_vars.items():
                if var_value is not None:
                    del var_value
                    deleted_vars.append(var_name)
            
            # Light garbage collection for text processing
            gc.collect()
            
            memory_after = psutil.Process().memory_info().rss / 1024 / 1024
            memory_freed = memory_before - memory_after
            
            if deleted_vars and memory_freed > 0:
                print(f"🧹 Text cleanup: deleted {', '.join(deleted_vars)}, freed {memory_freed:.1f} MB")
            
            return memory_freed
            
        except ImportError:
            # If psutil not available, just do garbage collection
            import gc
            gc.collect()
            return 0.0

if __name__ == "__main__":
    import sys
    
    processor = TextProcessor()
    
    if len(sys.argv) < 2:
        file_path = input("Text file path: ")
    else:
        file_path = sys.argv[1]
    
    if not os.path.exists(file_path):
        print(f"❌ File not found: {file_path}")
        sys.exit(1)
    
    result = processor.process_file(file_path)
    if result:
        print("\n" + "="*50)
        print("TEXT ANALYSIS RESULTS")
        print("="*50)
        print(json.dumps(result, indent=2)) 